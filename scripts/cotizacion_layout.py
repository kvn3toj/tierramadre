#!/usr/bin/env python3
"""
Ninguna caja de contenido puede cruzar el filete del pie ni salirse del margen.

Con el número de ítems variable, desbordar el pie deja de ser hipotético y pasa
a ser el fallo esperado: por eso esto corre dentro de la construcción y no como
herramienta suelta. Se mide sobre el .pptx, que es exacto, y no sobre el render.
"""
from pptx.util import Emu

PIE_Y, ANCHO, MARGEN = 1738, 1080, 88


def _px(v):
    return Emu(int(v)).inches * 96


def verifica(prs, saltar_primera=True):
    problemas = []
    for i, s in enumerate(prs.slides, 1):
        # la portada no lleva franja de pie: baja hasta el borde a propósito
        portada = saltar_primera and i == 1
        for sh in s.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            txt = sh.text_frame.text.replace("\n", " / ")[:40]
            arriba, alto = _px(sh.top), _px(sh.height)
            if not portada and arriba < PIE_Y and arriba + alto > PIE_Y - 1:
                problemas.append("lámina %d: cruza el pie (fin=%.0f) · %r"
                                 % (i, arriba + alto, txt))
            izq = _px(sh.left)
            # las bandas centradas a todo lo ancho (x=0, w=1080) son intencionales
            if izq >= MARGEN - 1 and izq + _px(sh.width) > ANCHO - MARGEN + 1:
                problemas.append("lámina %d: se sale del margen (der=%.0f) · %r"
                                 % (i, izq + _px(sh.width), txt))
    return problemas
