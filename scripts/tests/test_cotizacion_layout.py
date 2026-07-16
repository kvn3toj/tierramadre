from pptx import Presentation
from pptx.util import Emu, Pt


def _px(v):
    return Emu(int(round(v * 9525)))


def _lamina(prs, x, y, w, h, txt):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(_px(x), _px(y), _px(w), _px(h))
    tb.text_frame.text = txt
    return s


def _prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = _px(1080), _px(1920)
    return prs


def test_una_lamina_limpia_no_reporta_nada():
    from cotizacion_layout import verifica
    prs = _prs()
    _lamina(prs, 88, 1200, 500, 40, "Anillo Nexus")
    assert verifica(prs) == []


def test_detecta_texto_que_cruza_el_pie():
    from cotizacion_layout import verifica
    prs = _prs()
    prs.slides.add_slide(prs.slide_layouts[6])  # portada: la lámina 1 no lleva pie, así que la de prueba va en la 2
    _lamina(prs, 88, 1700, 500, 60, "se cuela bajo el pie")   # 1700+60 = 1760 > 1737
    problemas = verifica(prs)
    assert len(problemas) == 1
    assert "pie" in problemas[0]


def test_detecta_texto_que_se_sale_del_margen():
    from cotizacion_layout import verifica
    prs = _prs()
    _lamina(prs, 88, 1200, 950, 40, "demasiado ancho")        # 88+950 = 1038 > 992
    assert any("margen" in p for p in verifica(prs))


def test_la_banda_centrada_a_todo_lo_ancho_es_intencional():
    from cotizacion_layout import verifica
    prs = _prs()
    _lamina(prs, 0, 500, 1080, 30, "FOTOGRAFÍA EN PRODUCCIÓN")
    assert verifica(prs) == []


def test_la_portada_no_lleva_pie_y_puede_bajar_hasta_el_borde():
    # la portada no tiene franja de pie: sangra hasta el borde inferior a
    # propósito, así que una caja que "cruzaría" el pie en cualquier otra
    # lámina aquí debe ser ignorada por diseño, no reportada como falla
    from cotizacion_layout import verifica
    prs = _prs()
    _lamina(prs, 88, 1700, 500, 60, "portada sangrada")   # 1700+60 = 1760 > 1737
    assert verifica(prs) == []


def test_la_portada_sin_pie_igual_detecta_cajas_encimadas():
    """
    Un nombre de cliente largo desborda su caja y se monta sobre el precio
    total sin cruzar ningún pie (la portada no tiene) y sin salirse del
    margen (la caja sigue dentro de él) — el hueco exacto que dejaba pasar
    ese bug. saltar_primera evita el falso positivo del pie, no exime a la
    portada de comprobación entera: las cajas de texto no pueden encimarse.
    """
    from cotizacion_layout import verifica
    prs = _prs()
    s = _lamina(prs, 88, 1330, 900, 300, "María Fernanda Nombre Larguísimo")
    tb = s.shapes.add_textbox(_px(88), _px(1590), _px(500), _px(30))
    tb.text_frame.text = "Cotización N.° TM-2026-0043"
    problemas = verifica(prs)
    assert any("encima" in p for p in problemas)


def test_la_portada_con_cajas_apiladas_borde_con_borde_no_reporta_nada():
    """Tocarse borde con borde (como hoy: el titular termina justo donde
    empieza el subtítulo) es el diseño normal, no una falla."""
    from cotizacion_layout import verifica
    prs = _prs()
    s = _lamina(prs, 88, 1330, 900, 260, "Cotización\nSoul")
    tb = s.shapes.add_textbox(_px(88), _px(1590), _px(500), _px(30))  # 1330+260 = 1590: se tocan
    tb.text_frame.text = "Plan de producción · 240 unidades"
    assert verifica(prs) == []
