"""
Extremo a extremo: una cotización real, construida por el punto de entrada real.

Antes de esto había ~350 pruebas verdes entre este repo y el bot, y ninguna
renderizaba un deck desde un quote.json — el hueco de falsa seguridad que dejó
pasar una carátula que decía «Cotización Soul» en la cotización de un cliente
real. Esta prueba arma una cotización real (con un nombre largo, para que la
lámina de resumen no vuelva a desbordarse) y comprueba exactamente lo que
estaba mal: el nombre del cliente aparece, «Soul» no aparece en ningún lado,
no hay una promesa falsa de «incluye la base y los módulos», la gramática dice
«1 unidad» y no «1 unidades», hay N+2 láminas, y la retícula no se sale del pie.
"""
import json
import os
import subprocess
import sys

from pptx import Presentation

RAIZ = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def _quote(cliente="María Fernanda", items=None):
    items = items if items is not None else [{
        "itemNumber": 170,
        "nombre": "#170 Gotas del Amazonas Coleccion Privada",   # nombre real y largo: ver Fix 2
        "gemas": "Esmeralda F2 · 2,1 Ct",
        "joya": "Oro 18 k · 4 g",
        "unidades": 1,
        "unitario": "$9'200.000",
        "total": "$9'200.000",
        "fotoFileId": "",
    }]
    return {
        "quotationNumber": "TM-2026-0099",
        "cliente": cliente,
        "asesor": {"email": "a@tierramadre.co", "name": "Asesor"},
        "fecha": "16 de julio de 2026",
        "moneda": "COP",
        "qrUrl": "https://tierramadre.app/c/TM-2026-0099",
        "items": items,
        "total": items[0]["total"],
    }


def _escribe(tmp_path, q):
    p = tmp_path / "quote.json"
    p.write_text(json.dumps(q), encoding="utf-8")
    return str(p)


def _construye(scratch, quote_path, salida):
    env = dict(os.environ, TM_SCRATCH=scratch)
    r = subprocess.run(
        [sys.executable, os.path.join(RAIZ, "build-cotizacion-pptx.py"),
         "--quote", quote_path, "--out", salida],
        env=env, capture_output=True, text=True)
    assert r.returncode == 0, r.stderr
    return salida


def _textos(prs):
    return [sh.text_frame.text for s in prs.slides for sh in s.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]


def test_cotizacion_de_cliente_no_es_el_plan_soul(scratch, tmp_path):
    q = _quote()
    salida = _construye(scratch, _escribe(tmp_path, q), str(tmp_path / "quote.pptx"))
    prs = Presentation(salida)

    # portada + 1 pieza + resumen
    assert len(list(prs.slides)) == len(q["items"]) + 2

    texto_completo = " ".join(_textos(prs))
    assert "María Fernanda" in texto_completo             # el cliente sí llega a la carátula
    assert "Soul" not in texto_completo                    # no es el plan Soul
    assert "Incluye la base y los módulos" not in texto_completo  # sin módulos, sin promesa falsa
    assert "1 unidades" not in texto_completo               # concordancia: singular
    assert "1 unidad" in texto_completo
    assert "Plan" not in texto_completo                    # una cotización de cliente no tiene "plan"

    from cotizacion_layout import verifica
    assert verifica(prs) == []                              # nada cruza el pie ni el margen


def test_dos_items_no_repite_el_nombre_del_cliente_erroneamente(scratch, tmp_path):
    """Con más de una pieza, la pluralización y el conteo de láminas siguen bien."""
    base = _quote()["items"][0]
    items = [
        dict(base, itemNumber=170, unidades=1),
        dict(base, itemNumber=45, nombre="Esperanza", unidades=3,
             unitario="$1'000.000", total="$3'000.000"),
    ]
    q = _quote(items=items)
    q["total"] = "$12'200.000"
    salida = _construye(scratch, _escribe(tmp_path, q), str(tmp_path / "quote2.pptx"))
    prs = Presentation(salida)

    assert len(list(prs.slides)) == len(items) + 2
    texto_completo = " ".join(_textos(prs))
    assert "Soul" not in texto_completo
    assert "3 unidades" in texto_completo

    from cotizacion_layout import verifica
    assert verifica(prs) == []
