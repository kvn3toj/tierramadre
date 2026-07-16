#!/usr/bin/env python3
"""
Huella determinista de un .pptx: geometría, texto, color de relleno, color y
negrita de fuente, y hash de cada imagen.

Es la red del refactor: la lámina Soul no puede moverse ni un píxel — y un
píxel recoloreado también es un píxel movido (filetes, divisores y la barra
bicolor del pie viven de su color, no sólo de su posición). Se compara la
huella y no el render porque es exacta, no depende de LibreOffice ni de que
Cormorant esté instalada, y señala qué forma se movió (o de qué color) en vez
de «cambió el 3%».
"""
import hashlib

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def _px(v):
    """EMU -> px @96dpi, redondeado: el render trabaja en píxeles enteros"""
    return round(Emu(int(v)).inches * 96, 2)


def _rgb_hex(color):
    """RGBColor -> 'RRGGBB', o None si el color no es RGB explícito (heredado/tema)"""
    try:
        return str(color.rgb)
    except (AttributeError, TypeError, ValueError):
        return None


def huella(path):
    prs = Presentation(path)
    laminas = []
    for s in prs.slides:
        formas = []
        for sh in s.shapes:
            f = {
                "x": _px(sh.left), "y": _px(sh.top),
                "w": _px(sh.width), "h": _px(sh.height),
            }
            # el relleno sólido importa: filetes, divisores y la barra bicolor
            # del pie viven de su color, no sólo de su geometría. No toda forma
            # tiene relleno sólido (imágenes, formas sin relleno) — se protege
            # con try/except porque fill.type / fore_color.rgb lanzan en esos casos.
            try:
                if sh.fill.type == MSO_FILL.SOLID:
                    f["relleno"] = _rgb_hex(sh.fill.fore_color)
            except (AttributeError, TypeError, ValueError):
                pass
            if sh.has_text_frame and sh.text_frame.text.strip():
                f["texto"] = sh.text_frame.text
                runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
                if runs:
                    fu = runs[0].font
                    f["fuente"] = fu.name
                    f["tam"] = fu.size.pt if fu.size is not None else None
                    # bold puede ser None (heredado); el color de fuente puede
                    # ser heredado/tema y no tener .rgb — _rgb_hex ya lo cubre
                    f["negrita"] = fu.bold
                    f["color_fuente"] = _rgb_hex(fu.color)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # el byte de la imagen importa: el encuadre y el blanqueado viven ahí
                f["img"] = hashlib.sha256(sh.image.blob).hexdigest()[:16]
            formas.append(f)
        laminas.append(formas)
    return {"w": _px(prs.slide_width), "h": _px(prs.slide_height), "laminas": laminas}
