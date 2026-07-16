#!/usr/bin/env python3
"""
Cotización Soul → .pptx portrait 1080×1920 para subir a Drive (Google Slides).

Versión limpia: de 14 bloques por lámina a 6. Lo que se quitó y por qué:
  · «Cantidad» suelta      -> ya lo dice el eyebrow (Línea 04 · 10 unidades)
  · descripción larga      -> la repiten Gemas y Joya
  · «Por qué la recomendamos» -> una sola línea bajo la opción recomendada
  · nota al pie            -> detalle interno, no de cliente
  · sellos + QR grandes    -> son idénticos en las 8 láminas: pasan a franja de pie
La fotografía manda: ocupa el 55% superior.

Reutiliza los datos y las imágenes de build-cotizacion.py.

    TM_SCRATCH=<scratch> scripts/venv/bin/python scripts/build-cotizacion-pptx.py
"""
import io
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from importlib import import_module

base = import_module("build-cotizacion".replace("-", "_")) if False else None

# ── datos: se importan del generador HTML para no duplicar precios ───────────
import importlib.util
_spec = importlib.util.spec_from_file_location(
    "bc", os.path.join(os.path.dirname(os.path.abspath(__file__)), "build-cotizacion.py"))
bc = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(bc)

SCRATCH = os.environ.get("TM_SCRATCH", "")
OPT = os.path.join(SCRATCH, "opt")
ROOT = bc.ROOT
SALIDA = os.path.join(ROOT, "docs", "Cotizacion-Soul.pptx")

# Composites de ESTE render (QR, foto de portada, foto de cada pieza). TM_SCRATCH
# suele ser una carpeta compartida (el cotizador la fija a tests/fixtures por
# defecto) y dos renders concurrentes escribían el mismo qr.png / comp-portada.jpg:
# el QR del cliente A podía terminar en el deck del cliente B. main() la
# reemplaza por una subcarpeta única antes de dibujar nada; se limpia al salir.
# La caché de fotos (--fotos-cache) es la única excepción: esa sí es por-fileId
# y debe seguir siendo compartida entre renders.
COMPOSITES = SCRATCH

# ── lienzo ──────────────────────────────────────────────────────────────────
W, H = 1080, 1920                      # px
PX = Emu(9525)                         # 1 px @96dpi


def px(v):
    return Emu(int(round(v * 9525)))


def pt(v):                             # px -> Pt tipográfico
    return Pt(v * 0.75)


VINO = RGBColor(0x6E, 0x1D, 0x18)
VINO_SUAVE = RGBColor(0x5F, 0x2A, 0x22)
TINTA = RGBColor(0x3D, 0x20, 0x18)
TENUE = RGBColor(0x8A, 0x6A, 0x5C)
ORO = RGBColor(0xA8, 0x81, 0x3F)
VERDE = RGBColor(0x1F, 0x51, 0x30)
VERDE_HONDO = RGBColor(0x12, 0x33, 0x20)
PAPEL = RGBColor(0xFB, 0xF8, 0xF2)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
CREMA = RGBColor(0xFD, 0xF8, 0xF0)

SERIF = "Cormorant"
SANS = "Montserrat"

MARGEN = 88
FOTO_H = 1060
PIE_Y = 1738
BARRA_H = 20

# alto de línea real de Cormorant medido sobre el render, no el tamaño nominal
LINEA_CORMORANT = 1.15
# aire entre el eyebrow y el titular: eran 40 y sobraban. El hueco de arriba
# medía ~39 px y el de abajo 6; los 8 que se le quitan pagan el aire de la ficha
AIRE_EYEBROW = 32

# Filas del resumen. Soul tiene 5 y medían 100; el tope las deja intactas y el
# piso es lo que impide publicar una lámina con las filas encimadas.
ALTO_FILA_MIN = 72       # nombre 34 pt (~39 px a 1,15 em) + subtítulo 13 pt + aire
ALTO_FILA_MAX = 100      # el valor de Soul: no estirar una cotización de 2 ítems
ALTO_CIERRE_RESUMEN = 242   # filete + «Precio Total del Plan» + aclaración

# Retícula única de precios: las dos columnas viven SIEMPRE en la misma x,
# en todas las láminas, para que no salten al pasar de una a otra.
COL_U = W - MARGEN - 430          # precio por unidad
COL_T = W - MARGEN - 210          # precio total
FILA_H = 84                       # misma altura para recomendada y alternativas

# Los filetes iban en #F4F1EA sobre #FAF7F0: invisibles. Ahora se ven.
FILETE = RGBColor(0xD8, 0xC7, 0xB0)
FILETE_FINO = RGBColor(0xE6, 0xD9, 0xC6)


def encabeza_columnas(slide, y):
    texto(slide, COL_U, y, 200, 18, "Precio por Unidad", tam=11, color=TENUE,
          negrita=True, track=.18, mayus=True, alineado=PP_ALIGN.RIGHT)
    texto(slide, COL_T, y, 210, 18, "Precio Total", tam=11, color=TENUE,
          negrita=True, track=.18, mayus=True, alineado=PP_ALIGN.RIGHT)


def texto(slide, x, y, w, h, txt, *, fuente=SANS, tam=16, color=TINTA,
          negrita=False, track=0, mayus=False, alineado=PP_ALIGN.LEFT,
          interlinea=1.15, anclaje=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anclaje
    lineas = txt.split("\n")
    for i, ln in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alineado
        p.line_spacing = interlinea
        r = p.add_run()
        r.text = ln.upper() if mayus else ln
        f = r.font
        f.name = fuente
        f.size = pt(tam)
        f.bold = negrita
        f.color.rgb = color
        if track:
            # espaciado entre letras (no expuesto por python-pptx)
            r.font._rPr.set("spc", str(int(track * tam * 0.75 * 100)))
    return tb


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def img(slide, ruta, x, y, w=None, h=None):
    return slide.shapes.add_picture(ruta, px(x), px(y),
                                    px(w) if w else None, px(h) if h else None)


def _blanquea_fondo(im, tol=16):
    """
    Lleva el fondo del archivo a blanco puro.

    Los JPEG no traen el fondo en 255: el borde va de 242 a 255 (un velo de
    revelado). El aplanado de más abajo sólo tocaba de 250 en adelante, así que
    ese velo sobrevivía y dibujaba el rectángulo de cada foto contra el lienzo
    en 255: las costuras que se veían en la rejilla y en el trío.

    Se rellena sólo desde los bordes hacia dentro, no por umbral global: los
    brillos del oro y los diamantes quedan intactos porque la pieza los encierra.
    Con tol=16 no se come ni un píxel de pieza; de 28 en adelante el relleno se
    desborda, por eso el margen.
    """
    from PIL import ImageDraw
    im = im.copy()
    w, h = im.size
    for punto in ((0, 0), (w - 1, 0), (0, h - 1), (w - 1, h - 1),
                  (w // 2, 0), (w // 2, h - 1), (0, h // 2), (w - 1, h // 2)):
        ImageDraw.floodfill(im, punto, (255, 255, 255), thresh=tol)
    return im


def _abre(k):
    """abre una foto de pieza con el fondo ya llevado a blanco puro"""
    from PIL import Image
    return _blanquea_fondo(Image.open(ruta(k)).convert("RGB"))


def _encaja(im, w, h, cubrir=False, fondo=(255, 255, 255)):
    """devuelve im dentro de un lienzo w×h (contain, o cover recortando al centro)"""
    from PIL import Image
    iw, ih = im.size
    esc = max(w / iw, h / ih) if cubrir else min(w / iw, h / ih)
    im = im.resize((max(1, int(iw * esc)), max(1, int(ih * esc))), Image.LANCZOS)
    lienzo = Image.new("RGB", (int(w), int(h)), fondo)
    lienzo.paste(im, (int((w - im.width) / 2), int((h - im.height) / 2)))
    return lienzo


def compone_foto(p, destino):
    """
    Arma TODA el área de fotografía como una sola imagen 1080×FOTO_H.
    Hacerlo en PIL y no con varias imágenes en la lámina evita las costuras
    blancas donde dos fotos se tocan, y da control real del encuadre.
    """
    from PIL import Image
    f = p.get("foto_ruta") or p.get("foto")
    W_, H_ = int(W), int(FOTO_H)
    lienzo = Image.new("RGB", (W_, H_), (255, 255, 255))

    if f == "rejilla":
        cw, ch = W_ // 2, H_ // 2
        for i, k in enumerate(("mod-esmeralda", "mod-cornalina", "mod-citrino", "mod-zafiro")):
            celda = _encaja(_abre(k), cw, ch)
            lienzo.paste(celda, ((i % 2) * cw, (i // 2) * ch))
    elif f == "trio":
        # composición: el solitario manda a la izquierda, dos acompañan a la derecha
        gr = _encaja(_abre("anillo-1"), int(W_ * 0.58), int(H_ * 0.78))
        lienzo.paste(gr, (0, int(H_ * 0.11)))
        for i, k in enumerate(("anillo-2", "anillo-3")):
            ch = _encaja(_abre(k), int(W_ * 0.42), H_ // 2)
            lienzo.paste(ch, (int(W_ * 0.58), i * (H_ // 2)))
    elif p.get("sangra"):
        # a sangre el fondo del estudio es parte de la foto: blanquearlo la destruye
        lienzo = _encaja(Image.open(ruta(f)).convert("RGB"), W_, H_, cubrir=True)
    else:
        # margen de aire para que la pieza no toque los bordes
        lienzo = _encaja(_abre(f), int(W_ * 0.88), int(H_ * 0.88))
        fondo = Image.new("RGB", (W_, H_), (255, 255, 255))
        fondo.paste(lienzo, (int(W_ * 0.06), int(H_ * 0.06)))
        lienzo = fondo

    if not p.get("sangra") and f:
        # El JPEG deja el blanco de cada foto en ~253, y contra el lienzo en 255
        # se dibujaba el rectángulo de cada tesela. Aplanar lo casi-blanco a
        # blanco puro borra esas costuras sin tocar la pieza ni su sombra.
        import numpy as np
        a = np.asarray(lienzo).astype(np.int16)
        a[a.mean(axis=2) >= 250] = 255
        from PIL import Image as _I
        lienzo = _I.fromarray(a.astype("uint8"))

    lienzo.save(destino, "JPEG", quality=92, optimize=True, progressive=True)
    return destino


def alto_fila_resumen(n, banda):
    """
    Alto de fila del resumen para n ítems, o None si no cabe.

    El plan Soul trae 5 filas de 100. Con n variable la lámina se desbordaba por
    el pie, así que el alto sale de la banda disponible — pero con tope en 100
    para que Soul no se mueva, y con piso para no encimar el texto.
    """
    if n <= 0:
        return None
    alto = banda / float(n)
    if alto < ALTO_FILA_MIN:
        return None
    return min(alto, ALTO_FILA_MAX)


def _unidades_texto(n):
    """
    «1 unidad», no «1 unidades».

    Soul siempre trae 200/10/10/10/10 así que nunca se veía; una cotización de
    cliente con una sola pieza sí cae en el singular, y en un documento que le
    llega al cliente la concordancia importa.
    """
    try:
        una = int(n) == 1
    except (TypeError, ValueError):
        una = False
    return "1 unidad" if una else "%s unidades" % n


def alto_titulo(nombre, ancho=None, tam=62):
    """estima cuántas líneas ocupa el titular para que la ficha baje con él"""
    ancho = ancho or (W - MARGEN * 2)
    por_linea = max(1, int(ancho / (tam * 0.44)))
    lineas = 1
    largo = 0
    for palabra in nombre.split():
        n = len(palabra) + (1 if largo else 0)
        if largo + n > por_linea:
            lineas += 1
            largo = len(palabra)
        else:
            largo += n
    # Una línea no mide tam: mide el alto de línea de la fuente. Medido sobre el
    # render, Cormorant a 62 px con interlinea .98 avanza 70 px de base a base
    # (1,15 em). Devolver líneas × tam dejaba el titular de dos líneas 8 px más
    # corto de lo que ocupa y GEMAS se montaba sobre el descendente de «tipo
    # Cartier». El +16 es el aire bajo la última línea: sin él el hueco quedaba
    # en 6 px con dos líneas y en 14 px con una, y se notaba el salto.
    return lineas, int(tam * (1 + (lineas - 1) * LINEA_CORMORANT)) + 16


ANCHO_NOMBRE_RESUMEN = 500      # ancho de la caja de nombre en lamina_resumen


def _mide_nombre_resumen(nombre):
    """
    Alto de la caja de nombre en el resumen, y desplazamiento del subtítulo.

    lamina_pieza ya mide su titular con alto_titulo (línea 373) y por eso
    aguanta un nombre largo; lamina_resumen no lo hacía, tenía el alto de caja
    (36) y el offset del subtítulo (+46) fijos a mano, calibrados sólo para los
    nombres cortos de Soul. Con un nombre real y largo («#170 Gotas del Amazonas
    Coleccion Privada») el nombre se salía de una línea y el subtítulo quedaba
    atropellado, sin que cotizacion_layout.verifica pudiera detectarlo — mide
    la caja del cuadro de texto, no el texto ya envuelto dentro.

    A una sola línea (el caso de Soul, siempre) se devuelven los mismos números
    de antes tal cual: la huella dorada no puede moverse un píxel.
    """
    lineas, alto = alto_titulo(nombre, ancho=ANCHO_NOMBRE_RESUMEN, tam=34)
    if lineas == 1:
        return 36, 46
    return alto, alto + 4


def ruta(k, ext="jpg"):
    # una foto de cotización ya vive en la caché con su ruta absoluta
    if os.path.isabs(str(k)) and os.path.exists(str(k)):
        return k
    return os.path.join(OPT, "%s.%s" % (k, ext))


# Fotos de una cotización: llegan por fileId de Drive, no por clave de opt/.
# El plan Soul sigue resolviendo contra opt/ y no pasa por aquí.
FOTOS_CACHE = None
API_BASE = os.environ.get("TM_API_BASE", "https://tierramadre.app")


def resuelve_foto(p):
    """
    Deja p['foto'] como una ruta usable y decide el encuadre.

    Soul trae claves de opt/ ('brazalete'); una cotización trae fileId. Si la foto
    no baja, se vacía p['foto'] y la lámina cae al pergamino: una foto que falta
    no tumba la cotización.
    """
    from cotizacion_fotos import elige_encuadre, trae_foto
    from PIL import Image
    f = p.get("foto")
    if not f or FOTOS_CACHE is None:
        return p                       # camino Soul: intacto
    ruta_foto = trae_foto(f, FOTOS_CACHE, API_BASE)
    if not ruta_foto:
        p["foto"] = ""
        return p
    p["foto_ruta"] = ruta_foto
    p["sangra"] = elige_encuadre(Image.open(ruta_foto)) == "sangrar"
    return p


def qr_png(url, destino):
    import segno
    segno.make(url, error="m").save(destino, scale=10, border=0,
                                    dark="#6e1d18", light="#ffffff")
    return destino


def logo_claro(destino, color=(0xFD, 0xF8, 0xF0)):
    """
    Versión crema del logo para la portada.

    El archivo es tinta negra con alfa: sobre el crema del pie contrasta, pero
    sobre el degradado de la portada (luminancia ~31) quedaba en 1:1.3 y el
    descriptor «esmeraldas con ADN de paz» no se leía. Se recolorea la tinta y
    se conserva el alfa, así el trazo no se recorta.
    """
    from PIL import Image
    import numpy as np
    a = np.asarray(Image.open(ruta("logo", "png")).convert("RGBA")).copy()
    a[..., 0], a[..., 1], a[..., 2] = color
    Image.fromarray(a).save(destino)
    return destino


# ── franja de pie: logo · sellos · QR (misma en todas, queda de mobiliario) ──
def pie(slide, qr):
    rect(slide, 0, PIE_Y - 1, W, 1, RGBColor(0xE6, 0xDC, 0xCC))
    img(slide, ruta("logo", "png"), MARGEN, PIE_Y + 44, 232)
    sx = 470
    for i in (1, 2, 3):
        img(slide, os.path.join(ROOT, "public", "certification-logo-%d.png" % i),
            sx + (i - 1) * 76, PIE_Y + 30, 64)
    img(slide, qr, W - MARGEN - 104, PIE_Y + 26, 104)
    rect(slide, 0, H - BARRA_H, W * 0.62, BARRA_H, VERDE)
    rect(slide, W * 0.62, H - BARRA_H, W * 0.38, BARRA_H, VERDE_HONDO)


def foto_pieza(slide, p):
    if not p.get("foto"):
        # lámina de textura para la pieza sin fotografía
        comp = os.path.join(COMPOSITES, "comp-lamina.jpg")
        if not os.path.exists(comp):
            from PIL import Image
            _encaja(Image.open(ruta("pergamino")).convert("RGB"), int(W), int(FOTO_H),
                    cubrir=True).save(comp, "JPEG", quality=86)
        img(slide, comp, 0, 0, W, FOTO_H)
        img(slide, ruta("logo", "png"), W / 2 - 150, FOTO_H / 2 - 60, 300)
        texto(slide, 0, FOTO_H / 2 + 44, W, 30, "Fotografía en producción",
              tam=15, color=TENUE, negrita=True, track=.24, mayus=True,
              alineado=PP_ALIGN.CENTER)
    else:
        comp = os.path.join(COMPOSITES, "comp-%s.jpg" % p["key"])
        compone_foto(p, comp)
        img(slide, comp, 0, 0, W, FOTO_H)
    rect(slide, 0, FOTO_H - 2, W, 2, ORO)


def lamina_pieza(prs, p, qr):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, PAPEL)
    foto_pieza(s, p)

    y = FOTO_H + 48
    texto(s, MARGEN, y, 700, 24, "Línea %s · %s" % (p["linea"], _unidades_texto(p["unidades"])),
          tam=15, color=TENUE, negrita=True, track=.22, mayus=True)
    y += AIRE_EYEBROW
    nombre = p["nombre"].replace("\n", " ")
    _, alto = alto_titulo(nombre)          # el titular de 2 líneas empuja la ficha
    texto(s, MARGEN, y, W - MARGEN * 2, alto + 16, nombre,
          fuente=SERIF, tam=62, color=VINO, negrita=True, interlinea=0.98)
    y += alto + 26

    # dos líneas de ficha, sin filetes ni etiquetas pesadas
    for et, val in (("Gemas", p["gemas"]), ("Joya", p["joya"])):
        texto(s, MARGEN, y, 120, 20, et, tam=13, color=ORO, negrita=True, track=.2, mayus=True)
        texto(s, MARGEN + 130, y - 3, W - MARGEN * 2 - 130, 46, val, tam=17, color=TINTA,
              interlinea=1.3)
        y += 54

    y += 16
    if p.get("opciones"):
        # cabeceras: sin ellas el cliente no sabe cuál columna es unidad y cuál total
        encabeza_columnas(s, y)
        y += 26
        rect(s, MARGEN, y, W - MARGEN * 2, 1, FILETE)
        y += 18
        for o in p["opciones"]:
            rec = o.get("recomendada")
            # todas las filas miden lo mismo: el hueco del rótulo se reserva siempre
            if rec:
                rect(s, MARGEN, y - 6, 8, FILA_H - 6, ORO)
                texto(s, MARGEN + 22, y - 4, 300, 18, "Recomendada",
                      tam=12, color=ORO, negrita=True, track=.26, mayus=True)
            texto(s, MARGEN + 22, y + 22, 470, 40, o["nombre"],
                  fuente=SERIF, tam=34, color=VINO if rec else VINO_SUAVE, negrita=rec)
            texto(s, COL_U, y + 26, 200, 24, o["unitario"],
                  tam=19, color=VINO_SUAVE, negrita=True, alineado=PP_ALIGN.RIGHT)
            texto(s, COL_T, y + 26, 210, 24, o["total"],
                  tam=19, color=VINO, negrita=True, alineado=PP_ALIGN.RIGHT)
            y += FILA_H
            rect(s, MARGEN, y - 8, W - MARGEN * 2, 1, FILETE_FINO)
        if p.get("razon"):
            texto(s, MARGEN + 22, y + 12, W - MARGEN * 2 - 22, 48, p["razon"],
                  tam=15, color=VINO_SUAVE, interlinea=1.4)
    else:
        # misma retícula de columnas que las láminas con opciones
        encabeza_columnas(s, y)
        y += 26
        rect(s, MARGEN, y, W - MARGEN * 2, 2, VINO)
        y += 20
        texto(s, MARGEN + 22, y + 8, 470, 40, "Precio de la pieza",
              fuente=SERIF, tam=34, color=VINO_SUAVE)
        texto(s, COL_U, y + 12, 200, 30, p["unitario"],
              tam=22, color=VINO_SUAVE, negrita=True, alineado=PP_ALIGN.RIGHT)
        texto(s, COL_T, y + 12, 210, 30, p["total"],
              tam=22, color=VINO, negrita=True, alineado=PP_ALIGN.RIGHT)
    pie(s, qr)
    return s


def lamina_modulos(prs, m, qr):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, PAPEL)
    foto_pieza(s, {"foto": "rejilla", "key": "modulos"})

    y = FOTO_H + 48
    texto(s, MARGEN, y, 700, 24, "Línea %s · Composición" % m["linea"],
          tam=15, color=TENUE, negrita=True, track=.22, mayus=True)
    y += AIRE_EYEBROW
    # el titular reservaba 2 líneas a mano y sólo ocupa 1: los 40 px sobrantes
    # empujaban la aclaración por debajo del filete del pie (1737). Se mide igual
    # que en lamina_pieza para que la ficha baje sólo lo que el titular ocupa.
    nombre = m["nombre"].replace("\n", " ")
    _, alto = alto_titulo(nombre)
    texto(s, MARGEN, y, W - MARGEN * 2, alto + 16, nombre,
          fuente=SERIF, tam=62, color=VINO, negrita=True, interlinea=0.98)
    y += alto + 26
    texto(s, MARGEN, y, W - MARGEN * 2, 30, "Cada módulo se suma a la base del brazalete.",
          tam=17, color=TENUE)
    y += 46
    encabeza_columnas(s, y)
    y += 26
    rect(s, MARGEN, y, W - MARGEN * 2, 1, FILETE)
    y += 18
    for f in m["filas"]:
        texto(s, MARGEN + 22, y, 420, 34, "%s  %s" % (f["nombre"], f["cant"]),
              fuente=SERIF, tam=32, color=VINO_SUAVE, negrita=True)
        texto(s, MARGEN + 22, y + 40, 420, 20, f["sub"], tam=13, color=TENUE)
        texto(s, COL_U, y + 8, 200, 24, f["unitario"],
              tam=19, color=VINO_SUAVE, negrita=True, alineado=PP_ALIGN.RIGHT)
        texto(s, COL_T, y + 8, 210, 24, f["total"],
              tam=19, color=VINO, negrita=True, alineado=PP_ALIGN.RIGHT)
        y += 68
        rect(s, MARGEN, y - 10, W - MARGEN * 2, 1, FILETE_FINO)
    y += 14
    rect(s, MARGEN, y, W - MARGEN * 2, 2, VINO)
    y += 18
    texto(s, MARGEN + 22, y + 8, 460, 20, "Precio Total · módulos", tam=13, color=TENUE,
          negrita=True, track=.2, mayus=True)
    texto(s, COL_U, y, 430, 56, m["total"],
          tam=40, color=VINO, negrita=True, alineado=PP_ALIGN.RIGHT)
    # sin esto el cliente no sabe si los $49M se suman al plan o ya están dentro
    if m.get("aclaracion"):
        texto(s, MARGEN + 22, y + 62, W - MARGEN * 2 - 22, 24, m["aclaracion"],
              tam=15, color=VINO_SUAVE)
    pie(s, qr)
    return s


def velo(slide, x, y, w, h, color, alpha):
    """rectángulo con alfa real (python-pptx no expone transparencia)"""
    from pptx.oxml.ns import qn
    s = rect(slide, x, y, w, h, color)
    srgb = s.fill.fore_color._xFill.find(qn("a:srgbClr"))
    a = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
    srgb.append(a)
    return s


def lamina_portada(prs, qr, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    from PIL import Image
    comp = os.path.join(COMPOSITES, "comp-portada.jpg")
    base = _encaja(Image.open(ruta("portada")).convert("RGB"), int(W), int(H), cubrir=True)
    # degradado quemado en la imagen: dos rectángulos con alfa dejaban bandas duras
    capa = Image.new("RGB", base.size, (0x14, 0x0B, 0x07))
    mask = Image.new("L", (1, base.height))
    for yy in range(base.height):
        t = yy / base.height
        if t < 0.34:
            a = int(150 * (t / 0.34) ** 2)          # arriba apenas se insinúa
        else:
            a = int(150 + (245 - 150) * ((t - 0.34) / 0.66) ** 0.8)
        mask.putpixel((0, yy), min(a, 245))
    base = Image.composite(capa, base, mask.resize(base.size))
    base.save(comp, "JPEG", quality=88)
    img(s, comp, 0, 0, W, H)

    img(s, logo_claro(os.path.join(SCRATCH, "logo-claro.png")), MARGEN, 1120, 340)
    texto(s, MARGEN, 1290, 700, 24, d.EYEBROW_PORTADA,
          tam=15, color=RGBColor(0xDC, 0xB8, 0x72), negrita=True, track=.32, mayus=True)
    texto(s, MARGEN, 1330, W - MARGEN * 2, 260, d.TITULO_PORTADA,
          fuente=SERIF, tam=118, color=CREMA, negrita=True, interlinea=0.92)
    texto(s, MARGEN, 1590, 700, 30, d.SUBTITULO_PORTADA,
          tam=20, color=RGBColor(0xCF, 0xC3, 0xB2))
    rect(s, MARGEN, 1642, W - MARGEN * 2, 1, RGBColor(0x7A, 0x63, 0x40))
    texto(s, MARGEN, 1664, 500, 20, d.ETIQUETA_TOTAL,
          tam=13, color=RGBColor(0xDC, 0xB8, 0x72), negrita=True, track=.22, mayus=True)
    texto(s, MARGEN, 1692, 700, 90, d.TOTAL_PLAN, tam=64, color=CREMA, negrita=True)
    texto(s, MARGEN, 1800, 700, 24,
          "Valores en pesos colombianos (COP) · %s" % d.FECHA,
          tam=14, color=RGBColor(0x9A, 0x8C, 0x7C))
    img(s, qr, W - MARGEN - 116, 1690, 116)
    rect(s, 0, H - BARRA_H, W * 0.62, BARRA_H, VERDE)
    rect(s, W * 0.62, H - BARRA_H, W * 0.38, BARRA_H, VERDE_HONDO)
    return s


def lamina_resumen(prs, qr, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, PAPEL)
    y = 150
    texto(s, MARGEN, y, 700, 24, "Cierre", tam=15, color=TENUE,
          negrita=True, track=.22, mayus=True)
    y += 42
    texto(s, MARGEN, y, W - MARGEN * 2, 250, d.TITULO_RESUMEN,
          fuente=SERIF, tam=96, color=VINO, negrita=True, interlinea=0.94)
    y += 230
    texto(s, MARGEN, y, W - MARGEN * 2, 60, d.NOTA_RESUMEN,
          tam=18, color=TENUE, interlinea=1.4)
    y += 86
    encabeza_columnas(s, y)
    y += 26
    rect(s, MARGEN, y, W - MARGEN * 2, 1, FILETE)
    y += 22
    banda = PIE_Y - y - ALTO_CIERRE_RESUMEN
    alto_fila = alto_fila_resumen(len(d.RESUMEN), banda)
    if alto_fila is None:
        raise SystemExit("El resumen no cabe: %d ítems no entran sobre el pie. "
                         "Divide la cotización." % len(d.RESUMEN))
    for n, sub, u, pu, ptot in d.RESUMEN:
        alto_caja, offset_sub = _mide_nombre_resumen(n)
        texto(s, MARGEN + 22, y, ANCHO_NOMBRE_RESUMEN, alto_caja, n,
              fuente=SERIF, tam=34, color=VINO, negrita=True)
        # +46: a +40 los descendentes del nombre se comían el subtítulo
        texto(s, MARGEN + 22, y + offset_sub, 520, 22,
              "%s · %s" % (sub, _unidades_texto(u)), tam=13, color=VINO_SUAVE)
        texto(s, COL_U, y + 8, 200, 24, pu,
              tam=19, color=VINO_SUAVE, negrita=True, alineado=PP_ALIGN.RIGHT)
        texto(s, COL_T, y + 8, 210, 24, ptot,
              tam=19, color=VINO, negrita=True, alineado=PP_ALIGN.RIGHT)
        # nombre a una línea: la fila mide lo mismo de siempre (huella dorada
        # intacta). Nombre largo: la fila crece lo que haga falta para no
        # encimarse con la siguiente, sin tocar ALTO_FILA_MIN/MAX.
        y += alto_fila if alto_caja <= 36 else max(alto_fila, offset_sub + 40)
        rect(s, MARGEN, y - 18, W - MARGEN * 2, 1, FILETE_FINO)
    y += 26
    rect(s, MARGEN, y, W - MARGEN * 2, 2, VINO)
    y += 22
    texto(s, MARGEN + 22, y, 500, 20, d.ETIQUETA_TOTAL, tam=14, color=TENUE,
          negrita=True, track=.22, mayus=True)
    texto(s, MARGEN + 22, y - 14, W - MARGEN * 2 - 22, 120, d.TOTAL_PLAN,
          tam=76, color=VINO, negrita=True, alineado=PP_ALIGN.RIGHT)
    y += 132
    # Sólo el plan Soul trae módulos y alternativas que aclarar. Una cotización
    # de cliente sin ellos no repite la frase con datos ajenos: sin CIERRE_NOTA
    # no hay caja, y el silencio no miente sobre el alcance de lo cotizado.
    if d.CIERRE_NOTA:
        texto(s, MARGEN + 22, y, W - MARGEN * 2 - 22, 60, d.CIERRE_NOTA,
              tam=15, color=VINO_SUAVE, interlinea=1.5)
    pie(s, qr)
    return s


def main(argv=None):
    import argparse
    ap = argparse.ArgumentParser(description="Cotización Soul → .pptx 1080×1920")
    ap.add_argument("--out", default=SALIDA, help="ruta del .pptx de salida")
    ap.add_argument("--quote", help="quote.json; sin él se construye el plan Soul")
    ap.add_argument("--fotos-cache", default=os.path.join(SCRATCH, "fotos"),
                    help="caché de fotos bajadas de Drive")
    args = ap.parse_args(argv)
    if args.quote:
        from cotizacion_quote import carga_quote
        d = carga_quote(args.quote)
    else:
        d = bc
    if args.quote:
        globals()["FOTOS_CACHE"] = args.fotos_cache
        for p in d.PRODUCTOS:
            resuelve_foto(p)
    if not os.path.isdir(OPT):
        sys.exit("Faltan las fotos optimizadas en %s (define TM_SCRATCH)" % OPT)

    # Subcarpeta única para el QR y las fotos compuestas de ESTE render: dos
    # renders concurrentes sobre el mismo TM_SCRATCH ya no se pisan (ver
    # COMPOSITES arriba). La caché de fotos (FOTOS_CACHE) es aparte y sigue
    # compartida, por-fileId.
    import shutil
    import tempfile
    render_dir = tempfile.mkdtemp(prefix="render-", dir=SCRATCH)
    globals()["COMPOSITES"] = render_dir
    try:
        qr = qr_png(d.QR_URL, os.path.join(COMPOSITES, "qr.png"))

        prs = Presentation()
        prs.slide_width = px(W)
        prs.slide_height = px(H)

        lamina_portada(prs, qr, d)
        for p in d.PRODUCTOS:
            lamina_pieza(prs, p, qr)
            # la línea de módulos es del plan Soul; una cotización suelta no la tiene
            if getattr(d, "MODULOS", None) and p["key"] == "brazalete":
                lamina_modulos(prs, d.MODULOS, qr)
        lamina_resumen(prs, qr, d)

        from cotizacion_layout import verifica
        problemas = verifica(prs)
        if problemas:
            # antes de subir, no después: una lámina de cliente no sale con el texto
            # metido bajo el pie
            raise SystemExit("La lámina no cuadra:\n  " + "\n  ".join(problemas))

        os.makedirs(os.path.dirname(os.path.abspath(args.out)), exist_ok=True)
        prs.save(args.out)
        print("→ %s  (%.2f MB, %d láminas, %dx%d px)"
              % (args.out, os.path.getsize(args.out) / 1024 / 1024,
                 len(prs.slides.__iter__.__self__._sldIdLst), W, H))
    finally:
        # los composites ya viven dentro del .pptx (add_picture copia los bytes):
        # no hace falta conservarlos, y dejarlos crecería TM_SCRATCH sin límite.
        shutil.rmtree(render_dir, ignore_errors=True)


if __name__ == "__main__":
    main()
